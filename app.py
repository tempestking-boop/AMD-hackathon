import io, json, re, base64, time, tempfile
from pathlib import Path
from datetime import datetime

import streamlit as st
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="Executive Intelligence Platform",
    page_icon="📊",
    layout="wide",
)

st.markdown("""
<style>
  .main { background: #F7F8FC; }
  .stButton>button { background:#00C2B2; color:#0D1B4B; font-weight:700; border:none; border-radius:8px; }
  .stButton>button:hover { background:#009E92; color:white; }
  div[data-testid="stMetric"] { background:white; border-radius:10px; padding:12px; border-top:3px solid #00C2B2; }
  .block-container { padding-top: 2rem; }
</style>
""", unsafe_allow_html=True)

# ── Sidebar — model config ─────────────────────────────────────────────────────
with st.sidebar:
    st.title("⚙️ Configuration")
    base_url    = st.text_input("API Base URL", "http://localhost:8000/v1")
    api_key     = st.text_input("API Key", "abc-123", type="password")
    model_name  = st.text_input("Model", "Qwen2.5-7B")
    temperature = st.slider("Temperature", 0.0, 1.0, 0.0, 0.05)
    max_tokens  = st.number_input("Max Tokens", 500, 8000, 4000, 500)
    st.divider()
    st.caption("Upload a CSV or TXT, configure your local LLM, then generate.")

# ── Palette ────────────────────────────────────────────────────────────────────
C_CHARCOAL  = RGBColor(0x1C, 0x1C, 0x2E)
C_NAVY      = RGBColor(0x0D, 0x1B, 0x4B)
C_TEAL      = RGBColor(0x00, 0xC2, 0xB2)
C_GOLD      = RGBColor(0xF5, 0xA6, 0x23)
C_RED_SOFT  = RGBColor(0xE5, 0x3E, 0x3E)
C_GREEN     = RGBColor(0x1A, 0xB3, 0x73)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_OFFWHITE  = RGBColor(0xF7, 0xF8, 0xFC)
C_MUTED     = RGBColor(0x6B, 0x7A, 0x99)
C_DARK_TEXT = RGBColor(0x1C, 0x1C, 0x2E)
PALETTE     = ["#00C2B2", "#0D1B4B", "#F5A623", "#1AB373", "#E53E3E", "#7B61FF", "#FF6B35", "#00B4D8"]
CHART_BG    = "#F7F8FC"

# ══════════════════════════════════════════════════════════════════════════════
# DATA PROCESSING
# ══════════════════════════════════════════════════════════════════════════════

def process_csv(df):
    lines = [f"Shape: {len(df):,} rows × {len(df.columns)} columns",
             f"Columns: {', '.join(df.columns)}", ""]
    num_df = df.select_dtypes(include="number")
    if not num_df.empty:
        lines += ["=== NUMERIC STATISTICS ===", num_df.describe().round(2).to_string(), ""]
        for col in num_df.columns:
            pc = num_df[col].pct_change().mean()
            if abs(pc) > 0.03:
                lines.append(f"  Trend: {col} avg {'+' if pc>0 else ''}{pc*100:.1f}% per period")
        lines.append("")
    cat_df = df.select_dtypes(include=["object", "category"])
    if not cat_df.empty:
        lines.append("=== CATEGORICAL BREAKDOWN ===")
        for col in cat_df.columns[:6]:
            top = df[col].value_counts().head(5)
            lines.append(f"  {col}: " + ", ".join(f"{v}({c})" for v, c in top.items()))
    missing = df.isnull().sum()
    missing = missing[missing > 0]
    if not missing.empty:
        lines += ["", "=== MISSING DATA ==="]
        for col, cnt in missing.items():
            lines.append(f"  {col}: {cnt} missing ({cnt/len(df)*100:.1f}%)")
    lines += ["", "=== SAMPLE (first 5 rows) ===", df.head(5).to_string(index=False)]
    return "\n".join(lines)

def process_txt(text, max_chars=8000):
    text = text.strip()
    return text if len(text) <= max_chars else text[:max_chars] + "\n[...truncated]"

def extract_kpis(df):
    kpis = []
    num_cols = df.select_dtypes(include="number").columns.tolist()
    for col in num_cols[:8]:
        series = df[col].dropna()
        if len(series) < 2: continue
        total = series.sum(); mean = series.mean()
        mid = len(series) // 2
        first_h = series.iloc[:mid].mean() if mid > 0 else mean
        second_h = series.iloc[mid:].mean()
        pct_chg = ((second_h - first_h) / first_h * 100) if first_h != 0 else 0
        label = col.replace("_", " ").title()
        if total > 1_000_000:
            val_str = f"${total/1_000_000:.1f}M" if any(k in col.lower() for k in ["revenue","sales"]) else f"{total/1_000_000:.1f}M"
        elif total > 1_000:
            val_str = f"{total:,.0f}"
        else:
            val_str = f"{mean:.1f} avg"
        delta_str = f"{pct_chg:+.1f}%" if abs(pct_chg) > 1 else "Stable"
        sentiment = "positive" if pct_chg > 1 else ("negative" if pct_chg < -1 else "neutral")
        kpis.append({"label": label, "value": val_str, "delta": delta_str, "sentiment": sentiment})
    return kpis[:6]

# ══════════════════════════════════════════════════════════════════════════════
# CHARTS
# ══════════════════════════════════════════════════════════════════════════════

def _style_ax(ax, title=""):
    ax.set_facecolor(CHART_BG)
    ax.tick_params(colors="#6B7A99", labelsize=9)
    for spine in ax.spines.values(): spine.set_visible(False)
    ax.yaxis.grid(True, color="#E0E0E0", linewidth=0.7, linestyle="--")
    ax.set_axisbelow(True)
    if title:
        ax.set_title(title, fontsize=11, fontweight="bold", color="#0D1B4B", pad=10, loc="left")

def make_trend_chart(df):
    num_cols = df.select_dtypes(include="number").columns.tolist()
    if not num_cols: return None
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()
    series_cols = num_cols[:3]
    x_labels = df[cat_cols[0]].astype(str) if cat_cols else df.index.astype(str)
    if len(df) > 24:
        step = max(1, len(df) // 24); df = df.iloc[::step]
        x_labels = df[cat_cols[0]].astype(str) if cat_cols else df.index.astype(str)
    fig, ax = plt.subplots(figsize=(9, 3.8), facecolor=CHART_BG)
    for i, col in enumerate(series_cols):
        vals = pd.to_numeric(df[col], errors="coerce").fillna(0).values
        ax.plot(range(len(vals)), vals, color=PALETTE[i], linewidth=2.2, marker="o", markersize=4, label=col.replace("_"," ").title())
        if i == 0: ax.fill_between(range(len(vals)), vals, alpha=0.08, color=PALETTE[0])
    tick_step = max(1, len(x_labels) // 8)
    ax.set_xticks(range(0, len(x_labels), tick_step))
    ax.set_xticklabels(list(x_labels)[::tick_step], rotation=30, ha="right", fontsize=8)
    if len(series_cols) > 1: ax.legend(fontsize=8, framealpha=0.5)
    _style_ax(ax, "Performance Trend")
    plt.tight_layout()
    return fig

def make_category_chart(df):
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()
    num_cols = df.select_dtypes(include="number").columns.tolist()
    if not cat_cols or not num_cols: return None
    grp = df.groupby(cat_cols[0])[num_cols[0]].sum().nlargest(8).sort_values()
    if grp.empty: return None
    fig, ax = plt.subplots(figsize=(9, 3.8), facecolor=CHART_BG)
    colors = ["#00C2B2" if i == len(grp) - 1 else "#0D1B4B" for i in range(len(grp))]
    bars = ax.barh(grp.index.astype(str), grp.values, color=colors, height=0.55, edgecolor="none")
    for bar, val in zip(bars, grp.values):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2, f"{val:,.0f}", va="center", fontsize=8, color="#6B7A99")
    title = f"{num_cols[0].replace('_',' ').title()} by {cat_cols[0].replace('_',' ').title()}"
    _style_ax(ax, title)
    plt.tight_layout()
    return fig

def make_distribution_chart(df):
    cat_cols = df.select_dtypes(include=["object", "category"]).columns.tolist()
    num_cols = df.select_dtypes(include="number").columns.tolist()
    if not num_cols: return None
    if len(cat_cols) < 2:
        fig, ax = plt.subplots(figsize=(9, 3.8), facecolor=CHART_BG)
        ax.hist(df[num_cols[0]].dropna(), bins=12, color="#00C2B2", edgecolor="none", alpha=0.85)
        _style_ax(ax, f"Distribution — {num_cols[0].replace('_',' ').title()}")
        plt.tight_layout()
        return fig
    pivot = df.groupby([cat_cols[0], cat_cols[1]])[num_cols[0]].sum().unstack(fill_value=0)
    if pivot.empty or len(pivot.columns) > 8: return None
    fig, ax = plt.subplots(figsize=(9, 3.8), facecolor=CHART_BG)
    pivot.plot(kind="bar", ax=ax, color=PALETTE[:len(pivot.columns)], edgecolor="none", width=0.65)
    ax.legend(fontsize=8, framealpha=0.5)
    ax.set_xticklabels(ax.get_xticklabels(), rotation=30, ha="right", fontsize=8)
    _style_ax(ax, f"{num_cols[0].replace('_',' ').title()} by {cat_cols[0].title()} & {cat_cols[1].title()}")
    plt.tight_layout()
    return fig

# ══════════════════════════════════════════════════════════════════════════════
# LLM + INSIGHT PIPELINE
# ══════════════════════════════════════════════════════════════════════════════

STEP1_SYSTEM = """You are a senior management consultant at a top-tier strategy firm.
You write board-ready executive briefings: direct, evidence-based, commercially sharp.
Every bullet must contain an observation AND its business significance."""

STEP1_USER = """Analyse the data below and produce a rich executive briefing.

FORMATTING RULES:
- Each bullet = 1 complete thought: observation + significance + implication.
- Target 20–35 words per bullet.
- Use numbers and metrics wherever possible.
- Executive voice: direct, commercial, data-driven, action-oriented.

═══════════════════════════════
EXECUTIVE SUMMARY (80–150 words):

BUSINESS HEALTH RATIONALE (2 sentences):

KEY FINDINGS (6 items, F1–F6):

TREND ANALYSIS (5 items, T1–T5):

RISK ASSESSMENT (5 items, R1–R5):

OPPORTUNITIES (5 items, O1–O5):

STRATEGIC PRIORITIES (3 items, each: Title / Rationale / Expected Outcome):

EXECUTIVE TAKEAWAY (1 sentence, 20–30 words):

DATA:
---
{data_summary}
---"""

STEP2_SYSTEM = "You are a JSON formatter. Output ONLY the JSON object. No markdown. No preamble."
STEP2_USER = """Convert this briefing to JSON exactly matching this schema:
{{
  "executive_summary": "string",
  "business_health_rationale": "string",
  "key_findings": ["string x6"],
  "trends": ["string x5"],
  "risks": ["string x5"],
  "opportunities": ["string x5"],
  "strategic_priorities": [{{"title":"","rationale":"","outcome":""}} x3],
  "executive_takeaway": "string"
}}
Output ONLY the JSON. Start with {{ end with }}.
BRIEFING:
{briefing}"""

REPAIR_USER = """Fix this JSON to match the schema. Output ONLY valid JSON.
Required: executive_summary(str), business_health_rationale(str),
key_findings(6 strings), trends(5 strings), risks(5 strings),
opportunities(5 strings), strategic_priorities(3 dicts with title/rationale/outcome),
executive_takeaway(str).
Broken JSON: {broken}"""

SCHEMA = {
    "executive_summary": str, "business_health_rationale": str,
    "key_findings": list, "trends": list, "risks": list,
    "opportunities": list, "strategic_priorities": list, "executive_takeaway": str,
}

def validate(data):
    for k, t in SCHEMA.items():
        if k not in data: return False, f"Missing: {k}"
        if not isinstance(data[k], t): return False, f"Wrong type: {k}"
        if t == list and not data[k]: return False, f"Empty list: {k}"
    return True, ""

def extract_json(text):
    text = text.strip()
    for attempt in [text, re.sub(r"```(?:json)?\s*", "", text).replace("```","").strip()]:
        try: return json.loads(attempt)
        except: pass
    m = re.search(r"\{[\s\S]*\}", text)
    if m:
        try: return json.loads(m.group())
        except: pass
    return None

def llm_call(client, system, user, model, temperature, max_tokens):
    r = client.chat.completions.create(
        model=model,
        messages=[{"role":"system","content":system},{"role":"user","content":user}],
        temperature=temperature, max_tokens=max_tokens,
    )
    return r.choices[0].message.content.strip()

def generate_insights(client, data_summary, model, temperature, max_tokens, status):
    status.update(label="Step 1/2: Generating executive briefing…")
    briefing = llm_call(client, STEP1_SYSTEM, STEP1_USER.format(data_summary=data_summary), model, temperature, max_tokens)
    status.update(label="Step 2/2: Converting to structured JSON…")
    raw = llm_call(client, STEP2_SYSTEM, STEP2_USER.format(briefing=briefing), model, temperature, max_tokens)
    insights = extract_json(raw)
    for attempt in range(2):
        if insights:
            ok, err = validate(insights)
            if ok: break
        broken = json.dumps(insights) if insights else raw
        repaired = llm_call(client, STEP2_SYSTEM, REPAIR_USER.format(broken=broken[:4000]), model, temperature, max_tokens)
        insights = extract_json(repaired)
    if not insights: raise ValueError("All JSON extraction attempts failed.")
    ok, err = validate(insights)
    if not ok: raise ValueError(f"Schema invalid after repair: {err}")
    return insights

# ══════════════════════════════════════════════════════════════════════════════
# HEALTH SCORE
# ══════════════════════════════════════════════════════════════════════════════

def calculate_health_score(kpis, insights):
    score = 0
    _no_placeholder = lambda x: x and str(x).strip().lower() not in {"not specified","n/a","none","not available","not provided"}
    if kpis:
        pts = 50 / len(kpis)
        for k in kpis:
            s = str(k.get("sentiment","neutral")).lower()
            score += pts if s == "positive" else (pts * 0.5 if s == "neutral" else 0)
    else:
        score += 25
    opps  = [x for x in insights.get("opportunities", []) if _no_placeholder(x)]
    risks = [x for x in insights.get("risks", []) if _no_placeholder(x)]
    finds = [x for x in insights.get("key_findings", []) if _no_placeholder(x)]
    if len(opps) >= 3:  score += 10
    if len(risks) <= 2: score += 10
    if len(finds) >= 3: score += 10
    for key in ("key_findings","trends","risks","opportunities"):
        if [x for x in insights.get(key,[]) if x]: score += 5
    return max(0, min(100, round(score)))

# ══════════════════════════════════════════════════════════════════════════════
# PPTX BUILD (same logic as notebook, condensed)
# ══════════════════════════════════════════════════════════════════════════════

def _bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    if color is None: color = C_DARK_TEXT
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text).strip() if text else ""
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return txb

def _bullets(slide, items, l, t, w, h, size=14, color=None):
    if color is None: color = C_DARK_TEXT
    items = [str(x).strip() for x in (items or []) if x and str(x).strip()]
    if not items: items = ["No data available."]
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size); run.font.color.rgb = color; run.font.name = "Calibri"

def _section_divider(prs, number, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_NAVY)
    _tb(slide, number, 8.5, 0.8, 4.5, 5.5, size=200, bold=True, color=RGBColor(0x17,0x2A,0x6A), align=PP_ALIGN.RIGHT, font="Cambria")
    _oval(slide, 0.7, 3.1, 0.18, 0.18, C_TEAL)
    _tb(slide, title, 1.1, 2.7, 9.0, 1.2, size=40, bold=True, color=C_WHITE, font="Cambria")
    if subtitle:
        _tb(slide, subtitle, 1.1, 4.0, 9.0, 0.6, size=16, italic=True, color=C_TEAL)

def _health_gauge(slide, score, cx, cy):
    col = C_GREEN if score >= 70 else (C_GOLD if score >= 40 else C_RED_SOFT)
    _oval(slide, cx-1.1, cy-1.1, 2.2, 2.2, C_OFFWHITE)
    _oval(slide, cx-0.85, cy-0.85, 1.7, 1.7, col)
    _oval(slide, cx-0.62, cy-0.62, 1.24, 1.24, C_WHITE)
    _tb(slide, str(score), cx-0.5, cy-0.38, 1.0, 0.6, size=28, bold=True, color=col, align=PP_ALIGN.CENTER)
    _tb(slide, "/100", cx-0.45, cy+0.05, 0.9, 0.3, size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

def _kpi_card(slide, label, value, delta, sentiment, l, t, w=2.9, h=1.6):
    col = C_GREEN if sentiment=="positive" else (C_RED_SOFT if sentiment=="negative" else C_MUTED)
    arrow = "▲" if sentiment=="positive" else ("▼" if sentiment=="negative" else "—")
    band = C_TEAL if sentiment=="positive" else (C_RED_SOFT if sentiment=="negative" else C_GOLD)
    _rect(slide, l, t, w, h, C_WHITE)
    _rect(slide, l, t, w, 0.07, band)
    _tb(slide, label, l+0.15, t+0.12, w-0.2, 0.35, size=10, color=C_MUTED)
    _tb(slide, value, l+0.12, t+0.42, w-0.2, 0.62, size=26, bold=True, color=C_DARK_TEXT)
    _tb(slide, f"{arrow}  {delta}", l+0.12, t+1.15, w-0.2, 0.30, size=11, bold=True, color=col)

def build_deck(insights, kpis, chart_figs, filename):
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    ts = datetime.now().strftime("%B %d, %Y")
    score = calculate_health_score(kpis, insights)

    # ── Cover ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_CHARCOAL)
    _rect(slide, 10.2, 0, 3.13, 7.5, RGBColor(0x00,0x9E,0x92))
    _rect(slide, 10.2, 0, 0.25, 7.5, RGBColor(0x00,0x7A,0x72))
    _tb(slide, "EXECUTIVE INTELLIGENCE PLATFORM", 0.8, 1.5, 9.0, 0.5, size=11, bold=True, color=C_TEAL)
    _tb(slide, "Executive\nInsight Report", 0.8, 2.1, 9.0, 2.4, size=48, bold=True, color=C_WHITE, font="Cambria")
    _tb(slide, "AI-Powered Business Intelligence · Strategic Analysis · Executive Briefing", 0.8, 4.6, 9.0, 0.5, size=13, italic=True, color=C_MUTED)
    _rect(slide, 0.8, 5.5, 4.5, 0.03, C_TEAL)
    _tb(slide, f"Source: {filename}", 0.8, 5.6, 5.0, 0.35, size=11, color=C_MUTED)
    _tb(slide, ts, 0.8, 5.95, 5.0, 0.35, size=11, color=C_MUTED)

    # ── Section 1 ──
    _section_divider(prs, "01", "The Situation", "What does the data tell us?")

    # ── Exec Summary ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)
    _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
    _tb(slide, "Executive Summary", 0.55, 0.15, 9.0, 0.78, size=28, bold=True, color=C_WHITE, font="Cambria")
    _rect(slide, 0.5, 1.25, 8.6, 4.1, C_WHITE)
    _tb(slide, "THE SITUATION", 0.75, 1.4, 4.0, 0.35, size=9, bold=True, color=C_TEAL)
    _tb(slide, insights.get("executive_summary",""), 0.75, 1.75, 8.15, 3.45, size=14, color=C_DARK_TEXT)
    _rect(slide, 9.4, 1.25, 3.4, 4.1, C_WHITE)
    _tb(slide, "BUSINESS HEALTH", 9.55, 1.42, 3.1, 0.35, size=9, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    _health_gauge(slide, score, cx=11.1, cy=3.05)
    _tb(slide, insights.get("business_health_rationale","")[:200], 9.55, 4.25, 3.1, 0.9, size=10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)
    _rect(slide, 0.5, 5.5, 12.33, 1.5, C_NAVY)
    _tb(slide, "KEY TAKEAWAY", 0.8, 5.6, 2.5, 0.3, size=9, bold=True, color=C_TEAL)
    _tb(slide, insights.get("executive_takeaway",""), 0.8, 5.92, 11.8, 0.85, size=13, bold=True, color=C_WHITE, italic=True)

    # ── KPI Dashboard ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)
    _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
    _tb(slide, "KPI Dashboard", 0.55, 0.15, 8.0, 0.78, size=28, bold=True, color=C_WHITE, font="Cambria")
    if kpis:
        xs = [0.55, 4.75, 8.95]; ys = [1.25, 3.15]
        for idx, kpi in enumerate(kpis[:6]):
            _kpi_card(slide, kpi["label"], kpi["value"], kpi["delta"], kpi["sentiment"], xs[idx%3], ys[idx//3])
    else:
        _tb(slide, "KPI data not available for this file type.", 1.0, 2.5, 11.0, 1.0, size=16, color=C_MUTED, align=PP_ALIGN.CENTER)

    # ── Findings, Trends, Risks, Opportunities ──
    def findings_slide(prs, heading, items, icon, accent):
        items = [str(x).strip() for x in (items or []) if x and str(x).strip()][:7]
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, C_OFFWHITE)
        _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
        _tb(slide, heading, 0.55, 0.15, 9.5, 0.78, size=28, bold=True, color=C_WHITE, font="Cambria")
        _oval(slide, 0.55, 1.2, 0.75, 0.75, accent)
        _tb(slide, icon, 0.55, 1.25, 0.75, 0.62, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if len(items) <= 3:
            _rect(slide, 0.55, 2.2, 12.25, 4.9, C_WHITE)
            _bullets(slide, items, 0.9, 2.4, 11.8, 4.5, size=15)
        else:
            mid = (len(items)+1)//2
            _rect(slide, 0.55, 2.2, 5.9, 4.9, C_WHITE)
            _bullets(slide, items[:mid], 0.85, 2.4, 5.5, 4.5, size=14)
            _rect(slide, 6.9, 2.2, 5.9, 4.9, C_WHITE)
            _bullets(slide, items[mid:], 7.2, 2.4, 5.5, 4.5, size=14)

    findings_slide(prs, "Key Findings", insights.get("key_findings",[]), "🔍", C_TEAL)

    # Trends with chart
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)
    _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
    _tb(slide, "Trend Analysis", 0.55, 0.15, 7.0, 0.78, size=28, bold=True, color=C_WHITE, font="Cambria")
    _oval(slide, 0.55, 1.25, 0.7, 0.7, C_TEAL)
    _rect(slide, 0.55, 2.15, 6.1, 5.0, C_WHITE)
    _bullets(slide, insights.get("trends",[])[:4], 0.85, 2.35, 5.7, 4.6, size=14)
    if "trend" in chart_figs:
        buf = io.BytesIO(); chart_figs["trend"].savefig(buf, format="png", dpi=150, bbox_inches="tight"); buf.seek(0)
        slide.shapes.add_picture(buf, Inches(6.95), Inches(1.25), Inches(5.9), Inches(5.9))
    else:
        _rect(slide, 6.95, 1.25, 5.9, 5.9, C_WHITE)
        _tb(slide, "Chart not available", 7.2, 3.5, 5.3, 0.5, size=13, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)

    findings_slide(prs, "Risk Assessment", insights.get("risks",[]), "⚠", C_RED_SOFT)
    findings_slide(prs, "Opportunities", insights.get("opportunities",[]), "💡", C_GREEN)

    # ── Section 2 ──
    _section_divider(prs, "02", "The Path Forward", "What leadership must do next.")

    # ── Strategic Priorities ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_OFFWHITE)
    _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
    _tb(slide, "Strategic Priorities", 0.55, 0.15, 9.0, 0.78, size=28, bold=True, color=C_WHITE, font="Cambria")
    _tb(slide, "Leadership Action Plan", 9.5, 0.28, 3.4, 0.5, size=11, color=C_GOLD, align=PP_ALIGN.RIGHT)
    priorities = insights.get("strategic_priorities", [])[:3]
    card_colors = [C_TEAL, C_GOLD, C_GREEN]; xs = [0.55, 4.88, 9.2]
    for i, (p, x, col) in enumerate(zip(priorities, xs, card_colors)):
        if isinstance(p, dict):
            title_t = str(p.get("title","")).strip()[:60]
            rat = str(p.get("rationale","")).strip()
            out = str(p.get("outcome","")).strip()
        else:
            parts = str(p).split(".",1)
            title_t = parts[0].strip()[:60]; rat = parts[1].strip() if len(parts)>1 else str(p); out = "Measure at next leadership review."
        _rect(slide, x, 1.3, 3.6, 5.8, C_WHITE)
        _rect(slide, x, 1.3, 3.6, 1.1, col)
        _tb(slide, f"0{i+1}", x+0.12, 1.32, 0.95, 1.05, size=38, bold=True, color=C_WHITE, font="Cambria")
        _tb(slide, title_t.upper(), x+1.05, 1.38, 2.4, 1.0, size=11, bold=True, color=C_WHITE)
        _tb(slide, "WHY NOW", x+0.18, 2.55, 3.0, 0.28, size=8, bold=True, color=col)
        _tb(slide, rat, x+0.18, 2.85, 3.25, 2.1, size=11, color=C_DARK_TEXT)
        _tb(slide, "EXPECTED OUTCOME", x+0.18, 5.05, 3.0, 0.28, size=8, bold=True, color=col)
        _tb(slide, out, x+0.18, 5.35, 3.25, 0.6, size=10, italic=True, color=C_MUTED)

    # ── Takeaway ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_CHARCOAL)
    _tb(slide, "\u201C", 0.5, 0.2, 2.5, 2.5, size=160, color=RGBColor(0x2A,0x2A,0x4A), bold=True, font="Cambria")
    _tb(slide, "THE BOTTOM LINE", 1.0, 1.0, 11.0, 0.45, size=11, bold=True, color=C_TEAL)
    _tb(slide, insights.get("executive_takeaway",""), 1.0, 1.55, 11.3, 2.5, size=26, bold=True, color=C_WHITE, font="Cambria")
    _rect(slide, 1.0, 6.5, 6.0, 0.04, C_TEAL)

    # ── Appendix charts ──
    for key, label in [("category","Appendix A — Category Breakdown"),("distribution","Appendix B — Distribution Analysis")]:
        if key in chart_figs:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _bg(slide, C_OFFWHITE)
            _rect(slide, 0, 0, 13.33, 1.05, C_CHARCOAL)
            _tb(slide, label, 0.55, 0.15, 12.3, 0.78, size=22, bold=True, color=C_WHITE, font="Cambria")
            buf = io.BytesIO(); chart_figs[key].savefig(buf, format="png", dpi=150, bbox_inches="tight"); buf.seek(0)
            slide.shapes.add_picture(buf, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.9))

    # ── Close ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, C_CHARCOAL)
    _rect(slide, 10.5, 0, 2.83, 7.5, RGBColor(0x00,0x9E,0x92))
    _tb(slide, "Thank You", 1.0, 2.1, 9.0, 1.6, size=54, bold=True, color=C_WHITE, font="Cambria")
    _tb(slide, "This report was generated automatically using AI.", 1.0, 3.9, 9.0, 0.6, size=14, italic=True, color=C_MUTED)
    _tb(slide, f"Source: {filename}", 1.0, 4.6, 9.0, 0.4, size=12, color=C_MUTED)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf

# ══════════════════════════════════════════════════════════════════════════════
# MAIN UI
# ══════════════════════════════════════════════════════════════════════════════

st.title("📊 Executive Intelligence Platform")
st.caption("Upload a CSV or TXT file, configure your local LLM, and generate a board-ready PowerPoint in seconds.")

uploaded = st.file_uploader("Upload CSV or TXT", type=["csv","txt"])

if uploaded:
    ext = Path(uploaded.name).suffix.lower()

    # ── Parse file ──
    if ext == ".csv":
        df = pd.read_csv(uploaded)
        st.success(f"**{uploaded.name}** — {len(df):,} rows × {len(df.columns)} columns")
        st.dataframe(df.head(5), use_container_width=True)
        data_summary = process_csv(df)
        kpis = extract_kpis(df)
    else:
        text = uploaded.read().decode("utf-8", errors="replace")
        df = None; kpis = []
        data_summary = process_txt(text)
        st.success(f"**{uploaded.name}** — {len(data_summary):,} characters")
        with st.expander("Preview"): st.text(data_summary[:1000])

    # ── KPI preview ──
    if kpis:
        st.subheader("KPIs detected")
        cols = st.columns(min(len(kpis), 3))
        for i, kpi in enumerate(kpis):
            with cols[i % 3]:
                delta_color = "normal" if kpi["sentiment"]=="positive" else ("inverse" if kpi["sentiment"]=="negative" else "off")
                st.metric(kpi["label"], kpi["value"], kpi["delta"], delta_color=delta_color)

    # ── Charts preview ──
    chart_figs = {}
    if df is not None:
        st.subheader("Charts")
        c1, c2, c3 = st.columns(3)
        fig = make_trend_chart(df)
        if fig: chart_figs["trend"] = fig; c1.pyplot(fig, use_container_width=True); plt.close(fig)
        fig = make_category_chart(df)
        if fig: chart_figs["category"] = fig; c2.pyplot(fig, use_container_width=True); plt.close(fig)
        fig = make_distribution_chart(df)
        if fig: chart_figs["distribution"] = fig; c3.pyplot(fig, use_container_width=True); plt.close(fig)

    st.divider()

    if st.button("🚀 Generate Executive Report", use_container_width=True):
        try:
            client = OpenAI(base_url=base_url, api_key=api_key)
            # Quick connection test
            with st.spinner("Testing connection…"):
                ids = [m.id for m in client.models.list().data]
            if model_name not in ids:
                st.warning(f"Model '{model_name}' not found. Available: {ids}")

            with st.status("Generating insights…", expanded=True) as status:
                insights = generate_insights(client, data_summary, model_name, temperature, max_tokens, status)
                status.update(label="Building PowerPoint deck…")
                pptx_buf = build_deck(insights, kpis, chart_figs, uploaded.name)
                status.update(label="✅ Done!", state="complete")

            # ── Preview insights ──
            st.subheader("Executive Summary")
            st.info(insights.get("executive_summary",""))

            col1, col2 = st.columns(2)
            with col1:
                st.subheader("Key Findings")
                for f in insights.get("key_findings",[]):
                    st.markdown(f"▸ {f}")
            with col2:
                st.subheader("Opportunities")
                for o in insights.get("opportunities",[]):
                    st.markdown(f"▸ {o}")

            st.subheader("Strategic Priorities")
            pc = st.columns(3)
            for i, p in enumerate(insights.get("strategic_priorities",[])[:3]):
                with pc[i]:
                    title = p.get("title","") if isinstance(p,dict) else str(p)[:60]
                    rat = p.get("rationale","") if isinstance(p,dict) else ""
                    out = p.get("outcome","") if isinstance(p,dict) else ""
                    st.markdown(f"**{title}**")
                    st.caption(rat)
                    if out: st.markdown(f"*{out}*")

            st.divider()
            st.download_button(
                label="⬇️ Download Executive_Intelligence_Report.pptx",
                data=pptx_buf,
                file_name="Executive_Intelligence_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

        except Exception as e:
            st.error(f"Error: {e}")
            st.info("Make sure your local LLM server is running and the base URL / model name are correct.")
else:
    st.info("👆 Upload a CSV or TXT file to get started.")
